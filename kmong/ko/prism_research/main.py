import os
import subprocess
import pdfplumber
import pandas as pd
from pptx import Presentation
from docx import Document
import multiprocessing

# 원본 폴더와 저장할 폴더 경로
source_folder = r"D:\GIT\pythoncrawling\kmong\ko\prism_research\prism_data_report_전체"
target_folder = r"D:\GIT\pythoncrawling\kmong\ko\prism_research\prism_data_report_전체_마이닝"

# hwp5txt.exe의 경로 (가상환경 내)
hwp5txt_path = r"D:\GIT\pythoncrawling\venv\Scripts\hwp5txt.exe"

# 저장할 폴더가 없으면 생성
os.makedirs(target_folder, exist_ok=True)

def clean_korean_text(text):
    """ 불필요한 줄바꿈 및 공백 정리만 수행 (형태소 분석 제거) """
    return text.replace("\n", " ").replace("\r", " ").strip()

def extract_text_from_pdf(file_path):
    """ PDF 파일에서 텍스트 추출 """
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return clean_korean_text(text.strip())

def extract_text_from_hwp(file_path):
    """ HWP 파일에서 텍스트 추출 (hwp5txt 사용) """
    try:
        result = subprocess.run([hwp5txt_path, file_path], capture_output=True, text=True, encoding="utf-8-sig")
        return clean_korean_text(result.stdout.strip())
    except Exception as e:
        return f"Error extracting HWP: {str(e)}"

def extract_text_from_xlsx(file_path):
    """ XLSX 파일에서 텍스트 추출 """
    text = ""
    df = pd.read_excel(file_path, sheet_name=None)  # 모든 시트 읽기
    for sheet_name, data in df.items():
        text += f"=== 시트: {sheet_name} ===\n"
        text += data.to_string(index=False) + "\n\n"
    return clean_korean_text(text.strip())

def extract_text_from_pptx(file_path):
    """ PPTX 파일에서 텍스트 추출 """
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return clean_korean_text(text.strip())

def extract_text_from_docx(file_path):
    """ DOCX 파일에서 텍스트 추출 """
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        text = f"Error extracting DOCX: {str(e)}"
    return clean_korean_text(text.strip())

def process_file(file_path):
    """ 개별 파일을 처리하는 함수 (병렬 실행) """
    file_name, file_ext = os.path.splitext(os.path.basename(file_path))
    txt_output_path = os.path.join(target_folder, f"{file_name}.txt")

    extracted_text = ""

    try:
        if file_ext.lower() == ".pdf":
            extracted_text = extract_text_from_pdf(file_path)
        elif file_ext.lower() == ".hwp":
            extracted_text = extract_text_from_hwp(file_path)
        elif file_ext.lower() == ".xlsx":
            extracted_text = extract_text_from_xlsx(file_path)
        elif file_ext.lower() in [".ppt", ".pptx"]:
            extracted_text = extract_text_from_pptx(file_path)
        elif file_ext.lower() in [".doc", ".docx"]:
            extracted_text = extract_text_from_docx(file_path)

        if extracted_text:
            with open(txt_output_path, "w", encoding="utf-8-sig") as f:
                f.write(extracted_text)
            print(f"✅ 변환 완료: {file_path} → {txt_output_path}")
        else:
            print(f"⚠️ 텍스트 추출 실패: {file_path}")

    except Exception as e:
        print(f"❌ 에러 발생 ({file_path}): {str(e)}")

if __name__ == "__main__":
    # 파일 목록 가져오기
    file_list = []
    for root, _, files in os.walk(source_folder):
        for file in files:
            file_list.append(os.path.join(root, file))

    # 병렬 처리 실행 (CPU 코어 수만큼 프로세스 생성)
    num_workers = min(multiprocessing.cpu_count(), len(file_list))
    with multiprocessing.Pool(num_workers) as pool:
        pool.map(process_file, file_list)
